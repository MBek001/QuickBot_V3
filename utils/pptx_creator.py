"""
Professional PowerPoint presentation generator.

IMPROVED VERSION:
- Better slide content
- Fixed closing slide
- Enhanced text fitting
"""

import logging
import re
from pathlib import Path
from typing import List, Optional, Dict

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("⚠️ python-pptx not installed. PPTX creation disabled.")

# ============================================================================
# COLOR THEMES
# ============================================================================

THEMES = {
    "professional": {
        "primary": RGBColor(31, 78, 121),  # Deep Blue
        "secondary": RGBColor(68, 114, 196),  # Medium Blue
        "accent": RGBColor(255, 192, 0),  # Gold
        "text": RGBColor(51, 51, 51),  # Dark Gray
        "background": RGBColor(255, 255, 255),  # White
    },
    "modern": {
        "primary": RGBColor(0, 120, 215),  # Azure
        "secondary": RGBColor(0, 153, 188),  # Teal
        "accent": RGBColor(232, 17, 35),  # Red
        "text": RGBColor(32, 31, 30),  # Almost Black
        "background": RGBColor(243, 242, 241),  # Light Gray
    },
    "vibrant": {
        "primary": RGBColor(142, 68, 173),  # Purple
        "secondary": RGBColor(230, 126, 34),  # Orange
        "accent": RGBColor(26, 188, 156),  # Turquoise
        "text": RGBColor(44, 62, 80),  # Navy
        "background": RGBColor(255, 255, 255),  # White
    },
    "corporate": {
        "primary": RGBColor(44, 62, 80),  # Dark Blue
        "secondary": RGBColor(52, 73, 94),  # Slate
        "accent": RGBColor(52, 152, 219),  # Light Blue
        "text": RGBColor(52, 73, 94),  # Slate
        "background": RGBColor(236, 240, 241),  # Off White
    },
}


# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def _slug(text: str) -> str:
    """Create filename-safe slug from text."""
    text = re.sub(r"\s+", "_", text.strip())
    text = re.sub(r"[^\w\-_.]", "", text, flags=re.UNICODE)
    return text[:50] or "presentation"


# ============================================================================
# SLIDE CREATORS
# ============================================================================

def _create_title_slide(prs, title: str, subtitle: str, theme: Dict):
    """Create professional title slide with proper text fitting."""
    if not PPTX_AVAILABLE:
        return None

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = theme["background"]

    # 🧹 CLEAN TITLE - Remove all markdown symbols
    clean_title = title.replace("**", "").replace("*", "").replace("__", "")
    clean_title = clean_title.replace("AI Generated", "").replace("Created by AI", "")
    clean_title = clean_title.strip()

    # Limit title length
    if len(clean_title) > 80:
        clean_title = clean_title[:77] + "..."

    # Title box - adjusted for better fit
    left = Inches(0.5)
    top = Inches(2.0)
    width = Inches(9)
    height = Inches(2.0)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = clean_title
    title_frame.word_wrap = True
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    # Adjust font size based on title length
    if len(clean_title) > 60:
        p.font.size = Pt(36)
    elif len(clean_title) > 40:
        p.font.size = Pt(40)
    else:
        p.font.size = Pt(44)

    p.font.bold = True
    p.font.color.rgb = theme["primary"]

    # 🧹 CLEAN SUBTITLE
    if subtitle:
        clean_subtitle = subtitle.replace("**", "").replace("*", "").replace("__", "")
        clean_subtitle = clean_subtitle.replace("AI Generated", "").replace("Created by AI", "")
        clean_subtitle = clean_subtitle.strip()

        if not clean_subtitle or len(clean_subtitle) < 5:
            clean_subtitle = "Professional Presentation"

        # Limit subtitle length
        if len(clean_subtitle) > 100:
            clean_subtitle = clean_subtitle[:97] + "..."

        sub_top = Inches(4.5)
        sub_box = slide.shapes.add_textbox(left, sub_top, width, Inches(1.2))
        sub_frame = sub_box.text_frame
        sub_frame.text = clean_subtitle
        sub_frame.word_wrap = True

        p = sub_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.color.rgb = theme["secondary"]

    # Decorative accent bar
    accent_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(6.5), Inches(9), Inches(0.15)
    )
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = theme["accent"]
    accent_shape.line.fill.background()

    return slide


def _create_content_slide(
        prs,
        slide_title: str,
        bullets: List[str],
        theme: Dict,
        slide_num: int
):
    """Create content slide with LARGER FONTS and better spacing."""
    if not PPTX_AVAILABLE:
        return None

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = theme["background"]

    # Clean title
    clean_title = slide_title.replace("**", "").replace("*", "").replace("__", "")
    clean_title = clean_title.strip()

    # Limit title length
    if len(clean_title) > 70:
        clean_title = clean_title[:67] + "..."

    # Header bar
    header = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = theme["primary"]
    header.line.fill.background()

    # Title in header - LARGER FONT
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25), Inches(8.5), Inches(0.9)
    )
    title_frame = title_box.text_frame
    title_frame.text = clean_title
    title_frame.word_wrap = True
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = title_frame.paragraphs[0]

    # 🔧 FIXED: Larger title fonts
    if len(clean_title) > 50:
        p.font.size = Pt(28)  # Was 24
    elif len(clean_title) > 35:
        p.font.size = Pt(32)  # Was 28
    else:
        p.font.size = Pt(36)  # Was 32

    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Slide number
    num_box = slide.shapes.add_textbox(
        Inches(9), Inches(0.35), Inches(0.8), Inches(0.6)
    )
    num_frame = num_box.text_frame
    num_frame.text = str(slide_num)

    p = num_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Content area - MORE SPACE
    content_top = Inches(1.5)
    content_left = Inches(0.7)  # Was 0.8
    content_width = Inches(8.6)  # Was 8.4
    content_height = Inches(5.2)  # Was 5.0

    content_box = slide.shapes.add_textbox(
        content_left, content_top, content_width, content_height
    )
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    # Add bullets with LARGER FONT
    for i, bullet_text in enumerate(bullets):
        # Remove markdown
        clean_bullet = bullet_text.replace("**", "").replace("*", "").replace("__", "")
        clean_bullet = clean_bullet.strip()

        if not clean_bullet or len(clean_bullet) < 3:
            continue

        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = clean_bullet
        p.level = 0
        p.space_before = Pt(12) if i > 0 else Pt(0)  # More space between bullets

        # 🔧 FIXED: LARGER fonts for content
        if len(clean_bullet) > 120:
            p.font.size = Pt(16)  # Was 14
        elif len(clean_bullet) > 80:
            p.font.size = Pt(18)  # Was 15
        else:
            p.font.size = Pt(20)  # Was 16

        p.font.color.rgb = theme["text"]

    # Accent footer
    footer = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(7.3), Inches(10), Inches(0.2)
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = theme["accent"]
    footer.line.fill.background()

    return slide


def _create_closing_slide(prs, theme: Dict, custom_message: str = None):
    """Create professional closing slide."""
    if not PPTX_AVAILABLE:
        return None

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = theme["primary"]

    # Thank you text
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(1.5)

    thank_box = slide.shapes.add_textbox(left, top, width, height)
    thank_frame = thank_box.text_frame
    thank_frame.text = custom_message or "Thank You"

    p = thank_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    sub_box = slide.shapes.add_textbox(left, Inches(4.5), width, Inches(0.8))
    sub_frame = sub_box.text_frame
    sub_frame.text = "Questions & Discussion"

    p = sub_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.color.rgb = theme["accent"]

    return slide


# ============================================================================
# MAIN CREATION FUNCTION
# ============================================================================

def create_pptx(
        title: str,
        slides: List[str],
        output_dir: str = "/tmp",
        theme_name: str = "professional"
) -> Optional[str]:
    """
    Create premium-quality PowerPoint with proper text handling.

    🔧 FIXED: Now calls _create_closing_slide properly
    """
    if not PPTX_AVAILABLE:
        logger.error("❌ Cannot create PPTX: python-pptx not installed")
        return None

    try:
        logger.info(f"📊 Creating PPTX: theme={theme_name}, slides={len(slides)}")

        theme = THEMES.get(theme_name, THEMES["professional"])
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Parse first slide
        if slides:
            first_slide_lines = [
                ln.strip() for ln in slides[0].split("\n") if ln.strip()
            ]

            # Clean title
            main_title = first_slide_lines[0] if first_slide_lines else title
            main_title = main_title.replace("**", "").replace("*", "").replace("__", "").strip()

            # Clean subtitle
            subtitle = first_slide_lines[1] if len(first_slide_lines) > 1 else "Professional Presentation"
            subtitle = subtitle.replace("**", "").replace("*", "").replace("__", "").strip()
        else:
            main_title = title
            subtitle = "Professional Presentation"

        # Create title slide
        _create_title_slide(prs, main_title, subtitle, theme)

        # Create content slides
        for idx, content in enumerate(slides[1:], start=2):
            lines = [ln.strip() for ln in content.split("\n") if ln.strip()]
            if not lines:
                continue

            # Clean and extract data
            slide_title = lines[0][:70]  # Limit length
            bullets = [b for b in lines[1:] if len(b.strip()) > 3]

            # Limit bullets to fit on slide
            if len(bullets) > 6:
                bullets = bullets[:6]

            _create_content_slide(prs, slide_title, bullets, theme, idx)

        # 🔧 FIXED: Closing slide now properly called
        _create_closing_slide(prs, theme, "Thank You")
        logger.info("✅ Added closing slide")

        # Save
        Path(output_dir).mkdir(parents=True, exist_ok=True)
        fname = f"presentation_{_slug(main_title)}.pptx"
        path = str(Path(output_dir) / fname)

        prs.save(path)
        logger.info(f"✅ PPTX saved: {path}")

        return path

    except Exception as e:
        logger.error(f"❌ PPTX error: {e}", exc_info=True)
        return None
